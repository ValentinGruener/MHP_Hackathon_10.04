from pptx import Presentation
from pptx.util import Emu

from app.models import Severity


def _color_distance(hex1: str, hex2: str) -> int:
    """Manhattan distance in RGB space between two 6-char hex color strings."""
    r1, g1, b1 = int(hex1[0:2], 16), int(hex1[2:4], 16), int(hex1[4:6], 16)
    r2, g2, b2 = int(hex2[0:2], 16), int(hex2[2:4], 16), int(hex2[4:6], 16)
    return abs(r1 - r2) + abs(g1 - g2) + abs(b1 - b2)


def _is_color_allowed(rgb: str, palette: set, tolerance: int = 5) -> bool:
    """Return True if rgb is within tolerance of any color in palette."""
    return any(_color_distance(rgb, c) <= tolerance for c in palette)


def _size_in_range(size_pt: float, range_dict: dict) -> bool:
    """Return True if size_pt falls within the min/max range."""
    if not range_dict:
        return True
    return range_dict.get("min", 0) <= size_pt <= range_dict.get("max", 9999)


def _check_required_slides(prs, required_slides: list) -> list[dict]:
    """Check that required slide types exist at specified positions."""
    total = len(prs.slides)
    errors = []
    for req in required_slides:
        pos = req.get("position")
        slide_type = req.get("type", "unknown")
        if pos == "last":
            check_idx = total - 1
            pos_label = "last"
        elif isinstance(pos, int):
            check_idx = pos - 1  # 1-indexed → 0-indexed
            pos_label = str(pos)
        else:
            continue

        if check_idx < 0 or check_idx >= total:
            errors.append(_make_error(
                slide_number=max(1, check_idx + 1),
                error_type="missing_required_slide",
                severity=Severity.critical,
                description=f'Pflicht-Slide "{slide_type}" fehlt — erwartet an Position {pos_label}',
                suggestion=f'Slide vom Typ "{slide_type}" an Position {pos_label} einfügen',
                auto_fixable=False,
                shape=None,
            ))
            continue

        # Heuristic title match: check if slide title text loosely matches type
        slide = prs.slides[check_idx]
        title_shape = slide.shapes.title
        title_text = (title_shape.text.strip().lower() if title_shape and title_shape.has_text_frame else "")

        type_keywords = {
            "title":      ["title", "titel", "deckblatt", "cover"],
            "agenda":     ["agenda", "inhalt", "übersicht", "contents", "table of"],
            "disclaimer": ["disclaimer", "haftungsausschluss", "legal", "impressum"],
            "summary":    ["summary", "zusammenfassung", "fazit", "conclusion"],
        }
        keywords = type_keywords.get(slide_type.lower(), [slide_type.lower()])
        if not any(kw in title_text for kw in keywords) and title_text:
            errors.append(_make_error(
                slide_number=check_idx + 1,
                error_type="wrong_slide_type",
                severity=Severity.warning,
                description=f'Folie {check_idx + 1} sollte vom Typ "{slide_type}" sein — Titel lautet: "{title_shape.text.strip() if title_shape else "(kein Titel)"}"',
                suggestion=f'Titel auf einen "{slide_type}"-typischen Inhalt prüfen',
                auto_fixable=False,
                shape=None,
            ))

    return errors


def check_rules(pptx_path: str, rules: dict) -> tuple[list[dict], float]:
    """Run rule-based CD checks against a presentation."""
    prs = Presentation(pptx_path)
    errors = []
    total_elements = 0
    checked_elements = 0

    allowed_fonts = set(rules.get("allowed_fonts", []))
    color_palette = set(c.upper().lstrip("#") for c in rules.get("color_palette", []))
    color_tolerance = int(rules.get("color_tolerance", 5))
    logo_spec = rules.get("logo_position")
    logo_required_on = rules.get("logo_required_on", [])  # e.g. ["first", "last"]

    # Font size: support both legacy exact-values list and new range dicts
    allowed_sizes = set(rules.get("allowed_font_sizes", []))
    title_size_range = rules.get("title_size", {})  # {"min": 24, "max": 36}
    body_size_range = rules.get("body_size", {})    # {"min": 12, "max": 20}
    use_size_ranges = bool(title_size_range or body_size_range)

    # Slide count limits
    slide_limits = rules.get("slide_limits", {})
    total_slides = len(prs.slides)
    if slide_limits:
        min_slides = slide_limits.get("min")
        max_slides = slide_limits.get("max")
        if min_slides is not None and total_slides < min_slides:
            errors.append(_make_error(
                slide_number=1,
                error_type="too_few_slides",
                severity=Severity.warning,
                description=f"Präsentation hat nur {total_slides} Folie(n) — Minimum ist {min_slides}",
                suggestion=f"Mindestens {min_slides} Folien einplanen",
                auto_fixable=False,
                shape=None,
            ))
        if max_slides is not None and total_slides > max_slides:
            errors.append(_make_error(
                slide_number=total_slides,
                error_type="too_many_slides",
                severity=Severity.warning,
                description=f"Präsentation hat {total_slides} Folien — Maximum ist {max_slides}",
                suggestion=f"Auf maximal {max_slides} Folien kürzen",
                auto_fixable=False,
                shape=None,
            ))

    # Required slides structural check
    required_slides = rules.get("required_slides", [])
    if required_slides:
        errors.extend(_check_required_slides(prs, required_slides))

    # Logo required_on check
    logo_errors = _check_logo_required_on(prs, logo_required_on)
    errors.extend(logo_errors)

    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            total_elements += 1

            # Font & size checks
            if shape.has_text_frame:
                checked_elements += 1
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        # Font check
                        if run.font.name and allowed_fonts and run.font.name not in allowed_fonts:
                            errors.append(_make_error(
                                slide_number=slide_idx,
                                error_type="wrong_font",
                                severity=Severity.critical,
                                description=f'Schriftart "{run.font.name}" ist nicht im CD erlaubt',
                                suggestion=f'Erlaubte Schriften: {", ".join(sorted(allowed_fonts))}',
                                current_value=run.font.name,
                                expected_value=sorted(allowed_fonts)[0] if allowed_fonts else None,
                                auto_fixable=True,
                                shape=shape,
                            ))

                        # Font size check — range-based or exact-list
                        if run.font.size:
                            size_pt = run.font.size.pt
                            size_ok = True
                            if use_size_ranges:
                                # Use title range for title shapes, body range for others
                                is_title = (shape == slide.shapes.title)
                                rng = title_size_range if is_title else body_size_range
                                size_ok = _size_in_range(size_pt, rng)
                                if not size_ok and rng:
                                    errors.append(_make_error(
                                        slide_number=slide_idx,
                                        error_type="wrong_font_size",
                                        severity=Severity.warning,
                                        description=f"Schriftgröße {size_pt}pt außerhalb erlaubtem Bereich ({rng.get('min')}–{rng.get('max')}pt)",
                                        suggestion=f"Größe auf {rng.get('min')}–{rng.get('max')}pt anpassen",
                                        current_value=str(size_pt),
                                        expected_value=f"{rng.get('min')}-{rng.get('max')}",
                                        auto_fixable=True,
                                        shape=shape,
                                    ))
                            elif allowed_sizes:
                                size_ok = size_pt in allowed_sizes
                                if not size_ok:
                                    closest = min(allowed_sizes, key=lambda s: abs(s - size_pt))
                                    errors.append(_make_error(
                                        slide_number=slide_idx,
                                        error_type="wrong_font_size",
                                        severity=Severity.warning,
                                        description=f"Schriftgröße {size_pt}pt nicht im CD",
                                        suggestion=f"Nächste erlaubte Größe: {closest}pt",
                                        current_value=str(size_pt),
                                        expected_value=str(closest),
                                        auto_fixable=True,
                                        shape=shape,
                                    ))

                        # Color check — tolerance-based
                        try:
                            _rgb_val = run.font.color.rgb if run.font.color else None
                        except AttributeError:
                            _rgb_val = None  # _NoneColor / theme color has no .rgb
                        if _rgb_val and color_palette:
                            rgb = str(_rgb_val).upper()
                            if not _is_color_allowed(rgb, color_palette, color_tolerance):
                                closest = min(color_palette, key=lambda c: _color_distance(rgb, c))
                                errors.append(_make_error(
                                    slide_number=slide_idx,
                                    error_type="wrong_color",
                                    severity=Severity.warning,
                                    description=f"Farbe #{rgb} nicht in CD-Palette (Toleranz ±{color_tolerance})",
                                    suggestion=f"Nächste CD-Farbe: #{closest}",
                                    current_value=f"#{rgb}",
                                    expected_value=f"#{closest}",
                                    auto_fixable=True,
                                    shape=shape,
                                ))

            # Empty slide check (once per slide, after last shape)
            slide_text = "".join(
                s.text_frame.text
                for s in slide.shapes
                if s.has_text_frame
            ).strip()
            if not slide_text and shape == list(slide.shapes)[-1]:
                errors.append(_make_error(
                    slide_number=slide_idx,
                    error_type="empty_slide",
                    severity=Severity.critical,
                    description="Leere Folie ohne Textinhalt",
                    suggestion="Folie entfernen oder Inhalt hinzufügen",
                    auto_fixable=False,
                    shape=None,
                ))

        # Missing title check
        if not slide.shapes.title:
            errors.append(_make_error(
                slide_number=slide_idx,
                error_type="missing_title",
                severity=Severity.warning,
                description="Folie hat keinen Titel",
                suggestion="Titel-Platzhalter hinzufügen",
                auto_fixable=False,
                shape=None,
            ))

    coverage = (checked_elements / total_elements * 100) if total_elements > 0 else 100.0

    return errors, coverage


def _check_logo_required_on(prs, logo_required_on: list) -> list[dict]:
    """Check that logo images appear on required slides (first, last)."""
    if not logo_required_on:
        return []

    errors = []
    total = len(prs.slides)
    slide_map = {}
    for spec in logo_required_on:
        if spec == "first":
            slide_map[0] = "first"
        elif spec == "last":
            slide_map[total - 1] = "last"
        elif isinstance(spec, int):
            slide_map[spec - 1] = str(spec)

    for idx, label in slide_map.items():
        if idx < 0 or idx >= total:
            continue
        slide = prs.slides[idx]
        has_image = any(
            shape.shape_type == 13  # MSO_SHAPE_TYPE.PICTURE
            for shape in slide.shapes
        )
        if not has_image:
            errors.append(_make_error(
                slide_number=idx + 1,
                error_type="missing_logo",
                severity=Severity.critical,
                description=f"Logo fehlt auf {label} Folie (Folie {idx + 1})",
                suggestion="Unternehmenslogo gemäß CD-Vorgabe einfügen",
                auto_fixable=False,
                shape=None,
            ))

    return errors


def _make_error(
    slide_number: int,
    error_type: str,
    severity: Severity,
    description: str,
    suggestion: str | None = None,
    current_value: str | None = None,
    expected_value: str | None = None,
    auto_fixable: bool = False,
    shape=None,
) -> dict:
    pos = {}
    if shape is not None:
        try:
            pos = {
                "position_x": shape.left,
                "position_y": shape.top,
                "position_w": shape.width,
                "position_h": shape.height,
            }
        except Exception:
            pass

    return {
        "slide_number": slide_number,
        "engine": "rules",
        "error_type": error_type,
        "severity": severity.value,
        "description": description,
        "suggestion": suggestion,
        "current_value": current_value,
        "expected_value": expected_value,
        "auto_fixable": auto_fixable,
        **pos,
    }
