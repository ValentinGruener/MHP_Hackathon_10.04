import asyncio

import httpx
from pptx import Presentation

from app.config import settings
from app.models import Severity


async def check_languagetool(pptx_path: str, language: str = "de-DE") -> list[dict]:
    """Run LanguageTool spelling/grammar checks on all slides in parallel."""
    prs = Presentation(pptx_path)
    slides_text = []

    for slide_idx, slide in enumerate(prs.slides, start=1):
        text_parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_parts.append(shape.text_frame.text)
        full_text = "\n".join(text_parts).strip()
        if full_text:
            slides_text.append((slide_idx, full_text))

    if not slides_text:
        return []

    # Check all slides in parallel (batches of 10)
    errors = []
    batch_size = 10
    for i in range(0, len(slides_text), batch_size):
        batch = slides_text[i : i + batch_size]
        results = await asyncio.gather(
            *[_check_text(slide_idx, text, language) for slide_idx, text in batch],
            return_exceptions=True,
        )
        for result in results:
            if isinstance(result, Exception):
                continue
            errors.extend(result)

    return errors


async def _check_text(slide_number: int, text: str, language: str) -> list[dict]:
    """Send text to LanguageTool and convert matches to check results."""
    errors = []

    try:
        async with httpx.AsyncClient(timeout=30.0) as client:
            response = await client.post(
                settings.languagetool_url,
                data={
                    "text": text,
                    "language": language,
                    "enabledOnly": "false",
                },
            )
            response.raise_for_status()
            data = response.json()
    except (httpx.HTTPError, Exception):
        return []

    for match in data.get("matches", []):
        context = match.get("context", {})
        replacements = match.get("replacements", [])
        suggestion = replacements[0]["value"] if replacements else None

        severity = Severity.warning
        rule_category = match.get("rule", {}).get("category", {}).get("id", "")
        if rule_category in ("TYPOS", "SPELLER"):
            severity = Severity.critical

        # Extract the incorrect text
        offset = context.get("offset", 0)
        length = context.get("length", 0)
        context_text = context.get("text", "")
        incorrect_text = context_text[offset : offset + length] if context_text else ""

        errors.append({
            "slide_number": slide_number,
            "engine": "languagetool",
            "error_type": match.get("rule", {}).get("category", {}).get("id", "spelling"),
            "severity": severity.value,
            "description": match.get("message", "Sprachfehler gefunden"),
            "suggestion": suggestion,
            "current_value": incorrect_text,
            "expected_value": suggestion,
            "auto_fixable": suggestion is not None,
        })

    return errors
