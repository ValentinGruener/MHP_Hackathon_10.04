import asyncio
import json

import anthropic
from pptx import Presentation

from app.config import settings
from app.models import Severity

SYSTEM_PROMPT = """Du bist ein Prüfassistent für PowerPoint-Präsentationen.
Deine Aufgabe ist es, semantische Probleme in Folieninhalten zu finden:
- Platzhalter-Text ("Lorem ipsum", "Text hier einfügen", "Titel eingeben", etc.)
- Leere oder bedeutungslose Inhalte
- Inkonsistente Begriffe oder Abkürzungen
- Fehlende Standard-Folien (Agenda, Zusammenfassung)
- Unvollständige Inhalte (Aufzählungen mit nur einem Punkt, etc.)

Antworte NUR mit dem Tool-Aufruf. Keine zusätzliche Erklärung."""

CHECK_TOOL = {
    "name": "report_errors",
    "description": "Report semantic errors found in slides",
    "input_schema": {
        "type": "object",
        "properties": {
            "errors": {
                "type": "array",
                "items": {
                    "type": "object",
                    "properties": {
                        "slide_number": {"type": "integer"},
                        "error_type": {
                            "type": "string",
                            "enum": [
                                "placeholder_text",
                                "empty_content",
                                "inconsistent_terms",
                                "missing_standard_slide",
                                "incomplete_content",
                            ],
                        },
                        "severity": {
                            "type": "string",
                            "enum": ["critical", "warning", "info"],
                        },
                        "description": {"type": "string"},
                        "suggestion": {"type": "string"},
                        "current_value": {"type": "string"},
                    },
                    "required": ["slide_number", "error_type", "severity", "description"],
                },
            },
        },
        "required": ["errors"],
    },
}


async def check_haiku(pptx_path: str) -> list[dict]:
    """Run semantic checks via Claude Haiku 4.5 with tool_use."""
    if not settings.anthropic_api_key:
        return []

    prs = Presentation(pptx_path)

    # Extract text per slide
    slides_data = []
    for slide_idx, slide in enumerate(prs.slides, start=1):
        text_parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_parts.append(shape.text_frame.text)
        slides_data.append({
            "slide_number": slide_idx,
            "text": "\n".join(text_parts).strip(),
        })

    if not slides_data:
        return []

    # Batch slides
    batch_size = settings.haiku_max_slides_per_batch
    all_errors = []

    batches = [
        slides_data[i : i + batch_size]
        for i in range(0, len(slides_data), batch_size)
    ]

    results = await asyncio.gather(
        *[_check_batch(batch) for batch in batches],
        return_exceptions=True,
    )

    for result in results:
        if isinstance(result, Exception):
            continue
        all_errors.extend(result)

    return all_errors


async def _check_batch(slides_batch: list[dict]) -> list[dict]:
    """Send a batch of slides to Haiku for semantic checking."""
    client = anthropic.AsyncAnthropic(api_key=settings.anthropic_api_key)

    slides_text = "\n\n".join(
        f"--- Folie {s['slide_number']} ---\n{s['text'] or '(leer)'}"
        for s in slides_batch
    )

    try:
        response = await client.messages.create(
            model=settings.anthropic_model,
            max_tokens=2048,
            system=SYSTEM_PROMPT,
            tools=[CHECK_TOOL],
            tool_choice={"type": "tool", "name": "report_errors"},
            messages=[
                {
                    "role": "user",
                    "content": f"Prüfe diese Folien auf semantische Probleme:\n\n{slides_text}",
                },
            ],
        )
    except Exception:
        return []

    errors = []
    for block in response.content:
        if block.type == "tool_use" and block.name == "report_errors":
            for err in block.input.get("errors", []):
                errors.append({
                    "slide_number": err["slide_number"],
                    "engine": "haiku",
                    "error_type": err.get("error_type", "semantic"),
                    "severity": err.get("severity", "info"),
                    "description": err.get("description", ""),
                    "suggestion": err.get("suggestion"),
                    "current_value": err.get("current_value"),
                    "expected_value": None,
                    "auto_fixable": False,
                })

    return errors
