import shutil
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from sqlalchemy import select
from sqlalchemy.ext.asyncio import AsyncSession

from app.models import CheckResult, Correction, CorrectionStatus, Presentation as PresentationModel


# Types of corrections we can safely apply
SAFE_CORRECTIONS = {"wrong_font", "wrong_font_size", "wrong_color", "TYPOS", "SPELLER"}


async def apply_corrections(
    db: AsyncSession,
    presentation: PresentationModel,
    check_result_ids: list[int],
) -> list[dict]:
    """Apply selected corrections to a copy of the PPTX."""
    original_path = Path(presentation.original_pptx_path)
    corrected_path = original_path.with_stem(original_path.stem + "_corrected")

    # Always work on a copy
    shutil.copy2(original_path, corrected_path)

    # Load check results
    result = await db.execute(
        select(CheckResult).where(
            CheckResult.id.in_(check_result_ids),
            CheckResult.presentation_id == presentation.id,
            CheckResult.auto_fixable == True,
        )
    )
    check_results = result.scalars().all()

    if not check_results:
        return []

    # Group corrections by type for efficient application
    prs = Presentation(str(corrected_path))
    corrections = []

    for cr in check_results:
        correction_result = _apply_single_correction(prs, cr)
        correction = Correction(
            presentation_id=presentation.id,
            check_result_id=cr.id,
            before_value=cr.current_value,
            after_value=cr.expected_value or cr.suggestion,
            status=CorrectionStatus.applied if correction_result else CorrectionStatus.failed,
            applied_at=datetime.utcnow() if correction_result else None,
        )
        db.add(correction)
        corrections.append({
            "check_result_id": cr.id,
            "before": cr.current_value,
            "after": cr.expected_value or cr.suggestion,
            "status": "applied" if correction_result else "failed",
        })

    # Save corrected file
    prs.save(str(corrected_path))

    # Validate the output by re-parsing
    try:
        Presentation(str(corrected_path))
    except Exception:
        # Corrupted output, remove and report failure
        corrected_path.unlink(missing_ok=True)
        for c in corrections:
            c["status"] = "failed"
        return corrections

    presentation.corrected_pptx_path = str(corrected_path)
    await db.commit()

    return corrections


def _apply_single_correction(prs: Presentation, check_result: CheckResult) -> bool:
    """Apply a single correction to the presentation. Returns True on success."""
    try:
        slide_idx = check_result.slide_number - 1
        if slide_idx < 0 or slide_idx >= len(prs.slides):
            return False

        slide = prs.slides[slide_idx]

        if check_result.error_type == "wrong_font":
            return _fix_font(slide, check_result)
        elif check_result.error_type == "wrong_font_size":
            return _fix_font_size(slide, check_result)
        elif check_result.error_type == "wrong_color":
            return _fix_color(slide, check_result)
        elif check_result.error_type in ("TYPOS", "SPELLER"):
            return _fix_typo(slide, check_result)
        else:
            return False
    except Exception:
        return False


def _fix_font(slide, check_result: CheckResult) -> bool:
    """Replace wrong font with the expected font."""
    if not check_result.expected_value:
        return False

    fixed = False
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.font.name == check_result.current_value:
                    # Write explicit value, not theme-inherited
                    run.font.name = check_result.expected_value
                    fixed = True
    return fixed


def _fix_font_size(slide, check_result: CheckResult) -> bool:
    """Fix font size to the nearest allowed size."""
    if not check_result.expected_value:
        return False

    target_size = float(check_result.expected_value)
    current_size = float(check_result.current_value)
    fixed = False

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.font.size and abs(run.font.size.pt - current_size) < 0.1:
                    run.font.size = Pt(target_size)
                    fixed = True
    return fixed


def _fix_color(slide, check_result: CheckResult) -> bool:
    """This is intentionally conservative. Only fix if expected_value is provided."""
    if not check_result.expected_value:
        return False

    current_hex = check_result.current_value.lstrip("#").upper()
    target_hex = check_result.expected_value.lstrip("#").upper()

    if len(target_hex) != 6:
        return False

    fixed = False
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.font.color and run.font.color.rgb:
                    if str(run.font.color.rgb).upper() == current_hex:
                        run.font.color.rgb = RGBColor.from_string(target_hex)
                        fixed = True
    return fixed


def _fix_typo(slide, check_result: CheckResult) -> bool:
    """Replace a typo with the suggested correction."""
    if not check_result.current_value or not check_result.expected_value:
        return False

    fixed = False
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if check_result.current_value in run.text:
                    run.text = run.text.replace(
                        check_result.current_value,
                        check_result.expected_value,
                    )
                    fixed = True
    return fixed
