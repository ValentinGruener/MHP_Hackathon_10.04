import hashlib
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


def extract_cd_rules(pptx_path: Path) -> dict:
    """Extract Corporate Design rules from a master PPTX template."""
    prs = Presentation(str(pptx_path))

    fonts = set()
    font_sizes = set()
    colors = set()
    layouts = []

    # Extract from slide masters and layouts
    for slide_master in prs.slide_masters:
        _extract_from_shapes(slide_master.shapes, fonts, font_sizes, colors)

        for layout in slide_master.slide_layouts:
            layouts.append(layout.name)
            _extract_from_shapes(layout.placeholders, fonts, font_sizes, colors)

    # Extract from actual slides (sample content)
    for slide in prs.slides:
        _extract_from_shapes(slide.shapes, fonts, font_sizes, colors)

    # Extract theme colors
    theme_colors = _extract_theme_colors(prs)
    colors.update(theme_colors)

    # Detect logo (largest image on first slide or master)
    logo_spec = _detect_logo(prs)

    # Derive size ranges from observed font sizes
    title_sizes = [s for s in font_sizes if s >= 20]
    body_sizes  = [s for s in font_sizes if s < 20]
    title_size_range = (
        {"min": min(title_sizes), "max": max(title_sizes)} if title_sizes else {}
    )
    body_size_range = (
        {"min": min(body_sizes), "max": max(body_sizes)} if body_sizes else {}
    )

    rules = {
        # Legacy fields (kept for backwards compat)
        "allowed_fonts": sorted(fonts),
        "allowed_font_sizes": sorted(font_sizes),
        "color_palette": sorted(colors),
        "logo_position": logo_spec,
        "slide_layouts": layouts,
        "slide_width": prs.slide_width,
        "slide_height": prs.slide_height,
        # New JSON-schema fields (Miro board spec)
        "color_tolerance": 5,
        "title_size": title_size_range,
        "body_size": body_size_range,
        "logo_required_on": ["first", "last"],
        "required_slides": [],   # populated manually during onboarding
        "slide_limits": {},      # populated manually during onboarding
        "severity_weights": {"critical": 5, "warning": 2, "info": 0},
    }

    return rules


def _extract_from_shapes(shapes, fonts: set, font_sizes: set, colors: set):
    """Extract font, size, and color information from shapes."""
    for shape in shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                font = run.font
                if font.name:
                    fonts.add(font.name)
                if font.size:
                    fonts_pt = font.size.pt
                    font_sizes.add(fonts_pt)
                try:
                    if font.color and font.color.rgb:
                        colors.add(str(font.color.rgb))
                except AttributeError:
                    pass  # _NoneColor or theme colors have no .rgb


def _extract_theme_colors(prs: Presentation) -> set:
    """Extract colors from the presentation theme."""
    colors = set()
    try:
        theme = prs.slide_masters[0].element
        # Parse theme XML for color scheme
        nsmap = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
        for clr_elem in theme.iter("{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr"):
            val = clr_elem.get("val")
            if val:
                colors.add(val.upper())
    except (IndexError, AttributeError):
        pass
    return colors


def _detect_logo(prs: Presentation) -> dict | None:
    """Try to detect the logo image and its position."""
    # Check slide masters first, then first slide
    sources = list(prs.slide_masters)
    if prs.slides:
        sources.append(prs.slides[0])

    best_image = None
    for source in sources:
        for shape in source.shapes:
            if shape.shape_type == 13:  # Picture
                if best_image is None or (shape.width * shape.height < best_image["area"]):
                    # Prefer smaller images (logos tend to be smaller than full-page images)
                    best_image = {
                        "x": shape.left,
                        "y": shape.top,
                        "w": shape.width,
                        "h": shape.height,
                        "area": shape.width * shape.height,
                    }

    if best_image:
        best_image.pop("area")
        return best_image
    return None
